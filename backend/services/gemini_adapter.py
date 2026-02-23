"""Gemini adapter.
Provides functions used by the application:
 - generate_words_from_topic(topic) -> list of {"word", "clue"}
 - group_terms_with_ai(terms_with_context) -> list of {"term", "answer"}
 - generate_words_from_file(file_bytes, mime_type) -> list of {"word", "clue"}

Uses the `google-genai` SDK (the current, supported version).
Environment variables:
 - GEMINI_API_KEY: Your Gemini API key
 - GEMINI_MODEL: Model ID to use (default: gemini-1.5-flash)
"""
from typing import List, Dict
import os
import json
import tempfile

try:
    from google import genai
    from google.genai import types
except ImportError:
    genai = None
    types = None

MODEL_NAME = os.getenv("GEMINI_MODEL", "gemini-flash-latest")


def _get_client():
    if genai is None:
        raise RuntimeError("google-genai is not installed. Run: pip install google-genai")
    api_key = os.getenv("GEMINI_API_KEY") or os.getenv("VITE_GEMINI_API_KEY")
    if not api_key:
        raise RuntimeError("GEMINI_API_KEY environment variable is missing.")
    return genai.Client(api_key=api_key)


def _predict_text(prompt: str) -> str:
    client = _get_client()
    response = client.models.generate_content(
        model=MODEL_NAME,
        contents=prompt,
    )
    text = response.text
    # Strip any markdown code fences the model may add
    text = text.replace("```json", "").replace("```", "").strip()
    return text


def generate_words_from_topic(topic: str) -> List[Dict]:
    """Call Gemini to generate crossword words and clues for a topic.

    Returns list of {"word", "clue"}.
    """
    prompt = f"""
    Generate 10 crossword words and short clues about {topic}.
    Return exactly a JSON array of objects with the form:
    [{{"word": "...", "clue": "..."}}]
    The "word" must be a single continuous string of uppercase English letters (A-Z) with NO spaces or punctuation.
    The "clue" should be a clear, concise definition or hint.
    DO NOT wrap the response in markdown blocks like ```json. Output standard JSON array text ONLY.
    """
    try:
        text = _predict_text(prompt)
        parsed = json.loads(text)
        if isinstance(parsed, list):
            out = []
            for it in parsed:
                w = it.get("word") if isinstance(it, dict) else None
                c = it.get("clue") if isinstance(it, dict) else None
                if w:
                    out.append({"word": w.strip(), "clue": (c or "").strip()})
            return out
    except Exception as e:
        print(f"Gemini API error in generate_words_from_topic: {e}")
        print("Falling back to hardcoded computing clues...")

    return [
        {"word": "CPU", "clue": "The central processing unit, the brain of the computer"},
        {"word": "CODE", "clue": "Instructions written in a programming language"},
        {"word": "PYTHON", "clue": "A popular high-level programming language used for AI"},
        {"word": "GEMINI", "clue": "Google's powerful multimodal AI model"},
        {"word": "SOFTWARE", "clue": "Programs and other operating information used by a computer"},
    ]


def group_terms_with_ai(terms_with_context: List[Dict]) -> List[Dict]:
    """Ask Gemini to provide answers/definitions for extracted terms.

    Input: list of {"term", "slide", "context"}
    Output: list of {"term", "answer"}
    """
    lines = []
    for t in terms_with_context[:50]:
        term = t.get("term")
        ctx = t.get("context", "")
        lines.append(f"- {term} : {ctx}")

    prompt = f"""
    For each item below, provide a concise answer/definition. Return exactly a JSON array
    of objects with fields: {{"term": string, "answer": string}} and nothing else.
    DO NOT wrap the response in markdown blocks like ```json. Output standard JSON array text ONLY.
    Input:
    {chr(10).join(lines)}
    """
    try:
        text = _predict_text(prompt)
        parsed = json.loads(text)
        if isinstance(parsed, list):
            return parsed
    except Exception:
        pass

    out = []
    for t in terms_with_context[:20]:
        term = t.get("term") if isinstance(t, dict) else str(t)
        out.append({"term": term, "answer": f"Gemini mock answer for {term}"})
    return out


def generate_words_from_file(file_bytes: bytes, mime_type: str, file_name: str = "") -> List[Dict]:
    """Pass a file (image, pdf, pptx) directly to Gemini for multimodal extraction.

    For PPTX files, text is extracted locally (fast) and sent as a text prompt.
    For other files, the file is uploaded to the Gemini Files API.
    Returns list of {"word", "clue"}.
    """
    is_pptx = (
        file_name.lower().endswith(".pptx")
        or mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    # --- Fast path: extract PPTX text locally with python-pptx ---
    if is_pptx:
        try:
            import io
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            slide_texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

            combined_text = "\n".join(slide_texts[:200])  # cap at 200 text blocks

            prompt = f"""
            Below is text extracted from a presentation.
            Extract the main concepts and generate 10 crossword words and short clues about these concepts.
            Return exactly a JSON array of objects with the form:
            [{{"word": "...", "clue": "..."}}]
            The "word" must be a single continuous string of uppercase English letters (A-Z) with NO spaces or punctuation.
            The "clue" should be a clear, concise definition or hint based on the presentation content.
            DO NOT wrap the response in markdown blocks like ```json. Output standard JSON array text ONLY.

            Presentation text:
            {combined_text}
            """
            text = _predict_text(prompt)
            parsed = json.loads(text)
            if isinstance(parsed, list):
                out = []
                for it in parsed:
                    w = it.get("word") if isinstance(it, dict) else None
                    c = it.get("clue") if isinstance(it, dict) else None
                    if w:
                        out.append({"word": w.strip().upper(), "clue": (c or "").strip()})
                if out:
                    return out
        except Exception as e:
            print(f"PPTX text extraction error, falling back to file upload: {e}")

    # --- Slow path: upload file directly to Gemini Files API (images, PDFs, etc.) ---
    prompt = """
    Analyze the attached document or image.
    Extract the main concepts and generate 10 crossword words and short clues about these concepts.
    Return exactly a JSON array of objects with the form:
    [{"word": "...", "clue": "..."}]
    The "word" must be a single continuous string of uppercase English letters (A-Z) with NO spaces or punctuation.
    The "clue" should be a clear, concise definition or hint based on the document.
    DO NOT wrap the response in markdown blocks like ```json. Output standard JSON array text ONLY.
    """
    try:
        client = _get_client()
        ext = os.path.splitext(file_name)[1] if file_name else ""
        tmp_path = ""
        uploaded_file = None
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                tmp.write(file_bytes)
                tmp_path = tmp.name

            uploaded_file = client.files.upload(
                file=tmp_path,
                config=types.UploadFileConfig(mime_type=mime_type) if mime_type else None,
            )

            response = client.models.generate_content(
                model=MODEL_NAME,
                contents=[uploaded_file, prompt],
            )
        finally:
            if tmp_path and os.path.exists(tmp_path):
                os.remove(tmp_path)
            if uploaded_file:
                try:
                    client.files.delete(name=uploaded_file.name)
                except Exception:
                    pass

        text = response.text.replace("```json", "").replace("```", "").strip()
        parsed = json.loads(text)
        if isinstance(parsed, list):
            out = []
            for it in parsed:
                w = it.get("word") if isinstance(it, dict) else None
                c = it.get("clue") if isinstance(it, dict) else None
                if w:
                    out.append({"word": w.strip().upper(), "clue": (c or "").strip()})
            return out
    except Exception as e:
        print(f"Gemini API error for file upload: {e}")

    return [
        {"word": "MOCKFILE", "clue": f"Mock clue for {mime_type}"},
        {"word": "MULTIMODAL", "clue": "Another mock file clue"},
    ]