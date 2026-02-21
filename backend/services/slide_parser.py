"""Slide parser utilities.

This module extracts candidate terms/phrases from a PowerPoint (PPTX) file and
returns structured items suitable for passing to an AI for grouping or answering.

Main function:
 - extract_terms(file_bytes) -> list of {"term": str, "slide": int, "context": str}

Dependencies:
 - Requires `python-pptx` for real parsing. If it's not installed the function
   raises a clear RuntimeError with installation instructions.
"""

import io

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def extract_terms(file_bytes):
    # Development/mock mode: if python-pptx is not installed but the
    # SLIDE_PARSER_MOCK env var is set to a truthy value, return a small
    # mocked set of terms so the rest of the API works without the
    # dependency.
    import os

    if Presentation is None:
        if os.getenv("SLIDE_PARSER_MOCK", "false").lower() in ("1", "true", "yes"):
            # return a few example terms with minimal context
            return [
                {"term": "example_term", "slide": 1, "context": "Example slide content."},
                {"term": "sample_term", "slide": 2, "context": "Another slide content."},
            ]

        raise RuntimeError(
            "python-pptx is not installed. Install dependencies with:\n"
            "  python3 -m venv .venv && source .venv/bin/activate && pip install -r backend/requirements.txt\n"
            "Or install just python-pptx: pip install python-pptx\n"
            "If you want to run without installing python-pptx set SLIDE_PARSER_MOCK=true"
        )
    """Extract meaningful terms from a PPTX file and return structured data.

    Returns a list of objects: {"term": str, "slide": int, "context": str}
    This keeps simple heuristics (words and short phrases) but preserves the slide-level
    context so an AI can better determine answers/definitions.
    """

    presentation = Presentation(io.BytesIO(file_bytes))
    results = []

    for i, slide in enumerate(presentation.slides, start=1):
        slide_text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_parts.append(shape.text.strip())

        slide_text = "\n".join(slide_text_parts).strip()
        if not slide_text:
            continue

        # Heuristics: split by newlines to get candidate phrases and also words
        candidates = set()
        for part in slide_text.split('\n'):
            part = part.strip()
            if len(part) == 0:
                continue
            # If the part looks like a short phrase, keep it
            if 3 < len(part) <= 60 and ' ' in part:
                candidates.add(part.lower())
            else:
                # split into words for longer lines
                for w in part.split():
                    w = ''.join(ch for ch in w if ch.isalnum())
                    if len(w) > 3:
                        candidates.add(w.lower())

        for term in candidates:
            results.append({"term": term, "slide": i, "context": slide_text})

    # De-duplicate by term, keeping the first occurrence
    seen = set()
    dedup = []
    for r in results:
        if r["term"] not in seen:
            dedup.append(r)
            seen.add(r["term"])

    return dedup[:50]